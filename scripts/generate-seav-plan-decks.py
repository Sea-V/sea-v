#!/usr/bin/env python3
"""
Generate visual mind-map style PowerPoint decks for SEA-V planning documents.
Output: one .pptx per topic in ~/Desktop/SEA-V-Website-Plan/
"""

from __future__ import annotations

import math
import sys
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path.home() / "Desktop" / "SEA-V-Website-Plan"

# Maritime palette
NAVY = RGBColor(0x0A, 0x25, 0x40)
TEAL = RGBColor(0x1A, 0x8A, 0x8A)
SKY = RGBColor(0xE8, 0xF4, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0x4A, 0x55, 0x66)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
CORAL = RGBColor(0xD4, 0x6A, 0x4A)
FOREST = RGBColor(0x2D, 0x6A, 0x4F)
BRANCH_COLORS = [TEAL, RGBColor(0x2B, 0x6C, 0xB0), GOLD, CORAL, FOREST, RGBColor(0x7B, 0x5E, 0xA7)]


class Branch:
    def __init__(self, title: str, bullets: Iterable[str], color: RGBColor | None = None):
        self.title = title
        self.bullets = list(bullets)
        self.color = color


class DeckSpec:
    def __init__(
        self,
        rel_path: str,
        title: str,
        subtitle: str,
        center: str,
        branches: list[Branch],
        takeaway: str = "",
        layout: str = "radial",
    ):
        self.rel_path = rel_path
        self.title = title
        self.subtitle = subtitle
        self.center = center
        self.branches = branches
        self.takeaway = takeaway
        self.layout = layout


def set_slide_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, *, size=14, bold=False, color=SLATE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_round_box(slide, left, top, width, height, fill: RGBColor, line: RGBColor | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_title_slide(prs: Presentation, spec: DeckSpec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(2), Inches(0.4), "SEA-V", size=16, bold=True, color=TEAL)
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(8.4), Inches(1.2), spec.title, size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(0.8), Inches(3.4), Inches(8.4), Inches(0.8), spec.subtitle, size=20, color=SKY)
    add_textbox(slide, Inches(0.8), Inches(6.8), Inches(8), Inches(0.4), "Website Plan  ·  www.sea-v.com", size=11, color=RGBColor(0x99, 0xAA, 0xBB))


def add_branch_card(slide, x, y, branch: Branch, w=Inches(2.35), h=Inches(1.55)):
    color = branch.color or TEAL
    box = add_round_box(slide, x, y, w, h, WHITE, color)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    p = tf.paragraphs[0]
    p.text = branch.title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = color
    for bullet in branch.bullets[:3]:
        bp = tf.add_paragraph()
        bp.text = f"• {bullet}"
        bp.font.size = Pt(9)
        bp.font.color.rgb = SLATE
        bp.space_before = Pt(2)
    return box


def connect(slide, x1, y1, x2, y2, color: RGBColor):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(2)
    return conn


def add_radial_mindmap_slide(prs: Presentation, spec: DeckSpec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SKY)
    add_textbox(slide, Inches(0.5), Inches(0.25), Inches(9), Inches(0.4), spec.title, size=22, bold=True, color=NAVY)

    cx, cy = Inches(5.0), Inches(3.85)
    hub = add_round_box(slide, cx - Inches(1.1), cy - Inches(0.55), Inches(2.2), Inches(1.1), NAVY, TEAL)
    tf = hub.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = spec.center
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    n = len(spec.branches)
    radius_x = Inches(3.15)
    radius_y = Inches(2.05)
    card_w = Inches(2.35)
    card_h = Inches(1.55)

    for i, branch in enumerate(spec.branches):
        angle = (2 * math.pi * i / n) - math.pi / 2
        bx = cx + radius_x * math.cos(angle) - card_w / 2
        by = cy + radius_y * math.sin(angle) - card_h / 2
        card = add_branch_card(slide, bx, by, branch, card_w, card_h)
        color = branch.color or BRANCH_COLORS[i % len(BRANCH_COLORS)]
        hx = cx
        hy = cy
        tx = bx + card_w / 2
        ty = by + card_h / 2
        connect(slide, hx, hy, tx, ty, color)


def add_tree_slide(prs: Presentation, spec: DeckSpec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SKY)
    add_textbox(slide, Inches(0.5), Inches(0.25), Inches(9), Inches(0.4), spec.title, size=22, bold=True, color=NAVY)

    root = add_round_box(slide, Inches(3.8), Inches(0.85), Inches(2.4), Inches(0.65), NAVY, TEAL)
    root.text_frame.text = spec.center
    root.text_frame.paragraphs[0].font.size = Pt(13)
    root.text_frame.paragraphs[0].font.bold = True
    root.text_frame.paragraphs[0].font.color.rgb = WHITE
    root.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    cols = min(3, len(spec.branches))
    rows = math.ceil(len(spec.branches) / cols)
    start_y = Inches(1.85)
    gap_x = Inches(3.05)
    gap_y = Inches(1.75)
    start_x = Inches(0.55)

    for i, branch in enumerate(spec.branches):
        col = i % cols
        row = i // cols
        x = start_x + col * gap_x
        y = start_y + row * gap_y
        add_branch_card(slide, x, y, branch)
        connect(slide, Inches(5.0), Inches(1.5), x + Inches(1.175), y, branch.color or TEAL)


def add_takeaway_slide(prs: Presentation, spec: DeckSpec):
    if not spec.takeaway:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_textbox(slide, Inches(0.8), Inches(0.6), Inches(3), Inches(0.4), "Key takeaway", size=14, bold=True, color=TEAL)
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(8.4), Inches(2.5), spec.takeaway, size=28, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.8), Inches(6.5), Inches(8), Inches(0.5), spec.subtitle, size=14, color=SKY)


def build_deck(spec: DeckSpec) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, spec)
    if spec.layout == "tree":
        add_tree_slide(prs, spec)
    else:
        add_radial_mindmap_slide(prs, spec)
    add_takeaway_slide(prs, spec)

    out_dir = ROOT / Path(spec.rel_path).parent
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = ROOT / Path(spec.rel_path).with_suffix(".pptx")
    prs.save(str(out_path))
    return out_path


SPECS = [
    DeckSpec(
        "README.md",
        "SEA-V Master Plan",
        "Visual blueprint for the maritime career platform",
        "SEA-V Platform",
        [
            Branch("Vision", ["Digital career home", "For yacht crew", "Built for life at sea"], TEAL),
            Branch("Product", ["22 web pages", "15 app modules", "Public yacht CV"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Business", ["Free early access", "Organic growth", "Future Pro tier"], GOLD),
            Branch("Tech", ["Static HTML/JS", "Supabase backend", "Vercel hosting"], FOREST),
            Branch("Trust", ["Verified references", "MCA/PYA sea time", "Private payslips"], CORAL),
            Branch("Roadmap", ["Live today", "Hardening next", "Networking future"], RGBColor(0x7B, 0x5E, 0xA7)),
        ],
        "One platform where seafarers track, verify, and share their entire career.",
    ),
    DeckSpec(
        "01-vision-and-goals/what-sea-v-is.md",
        "What SEA-V Is",
        "The digital career home for yacht crew",
        "SEA-V",
        [
            Branch("Career vault", ["Sea time log", "Certificate store", "Vessel history"], TEAL),
            Branch("Professional CV", ["Modern yacht CV", "Not a PDF dump", "Structured data"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Verification", ["Signed references", "Expiry tracking", "Audit-ready records"], GOLD),
            Branch("Shareable", ["Public profile link", "Employer-ready view", "Crew controls privacy"], FOREST),
        ],
        "SEA-V replaces scattered paperwork with one secure, structured maritime profile.",
    ),
    DeckSpec(
        "01-vision-and-goals/problem-we-solve.md",
        "The Problem",
        "Why maritime careers need a better system",
        "Fragmented careers",
        [
            Branch("Paper chaos", ["Scattered docs", "Lost at crew change", "No single source"], CORAL),
            Branch("Outdated CVs", ["Static PDFs", "Hard to verify", "Always out of date"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Sea time gaps", ["Manual tallying", "MCA/PYA rules", "No audit trail"], TEAL),
            Branch("Certs expire", ["Miss renewals", "No reminders", "Files in email"], GOLD),
            Branch("References lost", ["Email threads", "No verification", "Slow hiring"], FOREST),
        ],
        "Crew lose time, trust, and opportunities because their career data lives everywhere except one place.",
    ),
    DeckSpec(
        "01-vision-and-goals/long-term-vision.md",
        "Long-Term Vision",
        "Building the trusted ecosystem for seafarers",
        "Trusted ecosystem",
        [
            Branch("Track", ["Full career timeline", "Every vessel & cert", "Always current"], TEAL),
            Branch("Verify", ["Identity checks", "Sea-time attestation", "Reference network"], GOLD),
            Branch("Share", ["Confident presentation", "Employer trust", "One link CV"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Connect", ["Crew networking", "Industry partners", "Training providers"], FOREST),
        ],
        "SEA-V becomes the industry-standard way maritime professionals prove who they are and what they've done.",
    ),
    DeckSpec(
        "01-vision-and-goals/target-audience.md",
        "Target Audience",
        "Who SEA-V is built for",
        "Yacht crew",
        [
            Branch("Deck", ["OOW & above", "Career progression", "Ticket building"], TEAL),
            Branch("Engineering", ["ETO / Chief", "Cert tracking", "Vessel systems"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Interior", ["Hospitality skills", "Specialist quals", "Guest-facing"], GOLD),
            Branch("Bridge team", ["Navigation log", "Passage records", "Miles sailed"], FOREST),
            Branch("Future buyers", ["Recruiters", "Captains hiring", "Agencies"], CORAL),
        ],
        "Primary user: the seafarer. Future customer: anyone who needs to trust their record.",
    ),
    DeckSpec(
        "02-business-model/current-model-free-early-access.md",
        "Current Business Model",
        "How SEA-V works today",
        "Free early access",
        [
            Branch("No paywall", ["Free for crew", "No Stripe yet", "Full feature access"], TEAL),
            Branch("Growth loop", ["Crew build profiles", "Share public CVs", "Employers discover SEA-V"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Partnerships", ["Contact page CTAs", "Feedback channel", "Early access list"], GOLD),
            Branch("Costs", ["Supabase hosting", "Vercel static", "Email via Edge Fn"], FOREST),
        ],
        "Win crew first with a free, indispensable career tool. Monetise once trust and volume exist.",
    ),
    DeckSpec(
        "02-business-model/future-monetization-ideas.md",
        "Future Revenue",
        "Strategic options not yet built",
        "Monetisation paths",
        [
            Branch("Pro tier", ["Premium exports", "Verification badges", "Priority support"], GOLD),
            Branch("Recruiter seats", ["Search verified crew", "Consent-based", "Agency accounts"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Partners", ["Training providers", "Cert sync", "Placement fees"], TEAL),
            Branch("Verification", ["Identity layer", "Sea-time attest", "Industry stamp"], FOREST),
        ],
        "Revenue follows trust: charge for verification, access, and tools that save employers time.",
    ),
    DeckSpec(
        "03-site-structure/sitemap-public.md",
        "Public Site Map",
        "Pages anyone can visit without logging in",
        "Public zone",
        [
            Branch("Entry", ["index.html Login", "signup.html Register"], TEAL),
            Branch("Marketing", ["about.html Story", "contact.html Support"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Legal", ["privacy.html", "terms.html"], SLATE),
            Branch("Shareable", ["public-profile.html", "Opt-in yacht CV"], GOLD),
            Branch("Verification", ["verify-reference.html", "Token link for referees"], FOREST),
        ],
        "8 public routes drive signup, trust, sharing, and reference verification.",
        layout="tree",
    ),
    DeckSpec(
        "03-site-structure/sitemap-app-modules.md",
        "App Module Map",
        "15 protected pages after login",
        "Logged-in app",
        [
            Branch("Hub", ["dashboard.html Overview"], NAVY),
            Branch("Identity", ["profile.html"], TEAL),
            Branch("Career record", ["vessels · seatime · navigation · tenders"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Skills", ["onboard experience · specialist quals"], GOLD),
            Branch("Compliance", ["certificates · references"], FOREST),
            Branch("Private", ["payslips (never public)"], CORAL),
            Branch("Engagement", ["milestones · hobbies"], RGBColor(0x7B, 0x5E, 0xA7)),
            Branch("Output", ["cv-generator.html Export"], TEAL),
        ],
        "Every module feeds the public yacht CV and CV generator.",
        layout="tree",
    ),
    DeckSpec(
        "07-roadmap/now-live.md",
        "Live Today",
        "What's shipped and working now",
        "Production",
        [
            Branch("Career vault", ["Profile · vessels · sea time", "Certificates · refs"], TEAL),
            Branch("Sharing", ["Public yacht CV", "Reference verify flow"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Engagement", ["Auto milestones", "Hex badges", "Trophy case"], GOLD),
            Branch("Export", ["CV generator", "CSV · ZIP · payslips"], FOREST),
            Branch("Navigation", ["Leaflet passage map", "Ports & miles"], CORAL),
        ],
        "SEA-V is a working product. Next focus: hardening, publish controls, recruiter value.",
    ),
    # --- 02 business model (continued) ---
    DeckSpec(
        "02-business-model/partnerships-and-ecosystem.md",
        "Partnerships & Ecosystem",
        "How SEA-V grows beyond individual crew",
        "Industry network",
        [
            Branch("Crew first", ["Free career vault", "Word of mouth", "Shared CV links"], TEAL),
            Branch("Employers", ["View public profiles", "Verify references", "Future recruiter seats"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Agencies", ["Placement partners", "Early access contact", "Co-branded profiles"], GOLD),
            Branch("Training", ["Certificate providers", "Catalog sync", "Expiry reminders"], FOREST),
            Branch("Feedback loop", ["contact.html", "User feedback modal", "Product iteration"], CORAL),
        ],
        "SEA-V wins by becoming the profile crew already have when industry partners show up.",
    ),
    # --- 03 site structure (continued) ---
    DeckSpec(
        "03-site-structure/user-journeys.md",
        "User Journeys",
        "How people move through SEA-V",
        "Three paths",
        [
            Branch("New crew", ["Sign up", "Build profile", "Log sea time"], TEAL),
            Branch("Active crew", ["Add certs & refs", "Earn milestones", "Export CV"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Sharing", ["Enable public CV", "Send link to employer", "Get hired"], GOLD),
            Branch("Referee", ["Email token link", "Review & sign", "No account needed"], FOREST),
            Branch("Employer view", ["Open public profile", "See verified data", "Never sees payslips"], CORAL),
        ],
        "Every journey ends with a clearer, more trusted career story.",
    ),
    DeckSpec(
        "03-site-structure/public-vs-private-data.md",
        "Public vs Private Data",
        "What leaves the account and what never does",
        "Privacy model",
        [
            Branch("Always private", ["Payslips", "Login credentials", "Draft CV cache"], CORAL),
            Branch("Opt-in public", ["public_enabled toggle", "Shareable link", "Per-profile gate"], TEAL),
            Branch("Can appear publicly", ["Vessels · sea time · certs", "Refs · milestones · hobbies"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Trust layer", ["Signed references", "RLS per user", "Signed file URLs"], FOREST),
            Branch("Future control", ["Per-section publish", "Record-level flags", "Hardening in progress"], GOLD),
        ],
        "Crew choose to share. Payslips and account data never appear on public profiles.",
    ),
    # --- 04 product features ---
    DeckSpec(
        "04-product-features/dashboard.md",
        "Dashboard",
        "Career overview hub after login",
        "Dashboard",
        [
            Branch("At a glance", ["Profile summary", "Recent activity", "Quick links"], TEAL),
            Branch("Career stats", ["Sea time totals", "Vessel count", "Cert status"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Map snippets", ["Navigation highlights", "Ports visited", "Passage overview"], FOREST),
            Branch("Milestones", ["Recent badges", "Progress hints", "Sidebar trophy case"], GOLD),
            Branch("Entry point", ["First page after login", "Routes to all modules", "Evaluates achievements"], CORAL),
        ],
        "The dashboard answers: where am I in my career right now?",
    ),
    DeckSpec(
        "04-product-features/profile-and-public-cv.md",
        "Profile & Public CV",
        "Identity plus the shareable yacht CV",
        "Profile",
        [
            Branch("Identity", ["Rank · nationality · bio", "Photo upload", "Availability"], TEAL),
            Branch("Travel docs", ["Passports & visas", "Salary expectations", "Contact details"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Public toggle", ["public_enabled switch", "Save then share", "Gate if disabled"], GOLD),
            Branch("Public CV", ["public-profile.html?p=", "Employer-ready layout", "CTA to sign up"], FOREST),
            Branch("Account", ["Password reset", "Delete account RPC", "Privacy controls"], CORAL),
        ],
        "Profile is the identity layer. Public CV is the product crew show the world.",
    ),
    DeckSpec(
        "04-product-features/vessels-and-seatime.md",
        "Vessels & Sea Time",
        "The core career record for ticket building",
        "Career record",
        [
            Branch("Vessels", ["Yacht specs & photos", "Date ranges", "Linked records"], TEAL),
            Branch("Sea time log", ["Actual · standby · yard", "Watchkeeping days", "MCA/PYA aligned"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Verification", ["Entry status flags", "File attachments", "Audit trail"], GOLD),
            Branch("Export", ["CSV download", "Feeds milestones", "Powers CV generator"], FOREST),
            Branch("Links", ["Vessel on each entry", "Navigation passages", "Auto badge triggers"], CORAL),
        ],
        "Sea time and vessels are the factual backbone every other module builds on.",
    ),
    DeckSpec(
        "04-product-features/navigation-map.md",
        "Navigation Map",
        "World passage chart and nautical miles",
        "Navigation",
        [
            Branch("Leaflet map", ["World chart view", "OSM / Carto tiles", "Port markers"], TEAL),
            Branch("Passages", ["Routes & waypoints", "Distance calc", "Linked to sea time"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Ports", ["Visit log", "Geographic history", "Career storytelling"], FOREST),
            Branch("Visual CV", ["Map on public profile", "Shows global experience", "Employer wow-factor"], GOLD),
        ],
        "Turn dry sea time into a visual story of where you've sailed.",
    ),
    DeckSpec(
        "04-product-features/certificates-and-compliance.md",
        "Certificates & Compliance",
        "STCW, medical, rank certs with expiry tracking",
        "Certificates",
        [
            Branch("Catalog", ["Dropdown templates", "Mandatory flags", "STCW / medical"], TEAL),
            Branch("Expiry", ["Status indicators", "Renewal visibility", "Never miss dates"], CORAL),
            Branch("Files", ["PDF uploads", "Private storage", "Signed URLs"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Export", ["ZIP download", "Accountant-ready pack", "Public profile section"], FOREST),
            Branch("Compliance", ["Ticket building aid", "Audit-ready records", "Rank progression"], GOLD),
        ],
        "Certificates prove qualification. SEA-V makes expiry visible before it's a problem.",
    ),
    DeckSpec(
        "04-product-features/references-verification.md",
        "References & Verification",
        "Captain and officer references with signed proof",
        "References",
        [
            Branch("Create ref", ["Draft in app", "Referee details", "Role & vessel link"], TEAL),
            Branch("Send verify link", ["Email token URL", "verify-reference.html", "No login for referee"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Referee signs", ["Review text", "Confirm or decline", "Signature captured"], GOLD),
            Branch("Verified badge", ["Shows on public CV", "Trust for employers", "Edge Fn + RPC"], FOREST),
            Branch("Resend / track", ["Status workflow", "Draft → sent → verified", "Resend option"], CORAL),
        ],
        "Verified references turn 'trust me' into 'here's proof'.",
    ),
    DeckSpec(
        "04-product-features/milestones-badges.md",
        "Milestones & Badges",
        "Gamified career achievements (UI: Milestones)",
        "Milestones",
        [
            Branch("Auto awards", ["Sea days thresholds", "Vessel types", "Watchkeeping firsts"], TEAL),
            Branch("Manual log", ["Choose milestone", "Add date & vessel", "Personal records"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Hex badges", ["Typography SVG art", "Page-coloured rings", "Trophy case UI"], GOLD),
            Branch("Progress", ["Next milestone hint", "Locked vs earned", "Recent activity feed"], FOREST),
            Branch("Engine", ["achievements-engine.js", "Runs on page load", "Dedupes auto rows"], CORAL),
        ],
        "Milestones make career progress visible and rewarding — like Strava for sea time.",
    ),
    DeckSpec(
        "04-product-features/cv-generator.md",
        "CV Generator",
        "Live data into a polished yacht CV for export",
        "CV Generator",
        [
            Branch("Live pull", ["Reads all modules", "No duplicate entry", "Always current"], TEAL),
            Branch("Layout", ["Professional sections", "Photo & branding", "Print-ready"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Customise", ["Section toggles", "Draft in localStorage", "Preview before export"], GOLD),
            Branch("Output", ["Print / save PDF", "Share same data as public CV", "Employer handoff"], FOREST),
        ],
        "One click from career vault to a CV that actually reflects today's record.",
    ),
    DeckSpec(
        "04-product-features/payslips-private.md",
        "Payslips (Private)",
        "Financial records that never go public",
        "Payslips",
        [
            Branch("Always private", ["Never on public CV", "Excluded from sharing", "RLS protected"], CORAL),
            Branch("Organise", ["UK tax year filter", "Month grouping", "PDF uploads"], TEAL),
            Branch("Export", ["Accountant pack", "ZIP download", "Year-end ready"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Use case", ["Tax returns", "Salary history", "Personal finance only"], FOREST),
        ],
        "Career data is shareable. Financial data stays locked to the account owner.",
    ),
    # --- 05 technical architecture ---
    DeckSpec(
        "05-technical-architecture/tech-stack.md",
        "Tech Stack",
        "What SEA-V is built with",
        "Architecture",
        [
            Branch("Frontend", ["Vanilla HTML/CSS/JS", "Multi-page app", "No bundler"], TEAL),
            Branch("Backend", ["Supabase Postgres", "Auth + Storage", "RPC + Edge Fn"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Maps", ["Leaflet 1.9.4", "OSM tiles", "navigation.html"], FOREST),
            Branch("Deploy", ["Vercel static", "Netlify config too", "CSP in vercel.json"], GOLD),
            Branch("Dev/test", ["python http.server", "test-site.mjs", "test-supabase.mjs"], CORAL),
        ],
        "Simple static frontend, serious backend — fast to ship, cheap to host.",
        layout="tree",
    ),
    DeckSpec(
        "05-technical-architecture/supabase-data-model.md",
        "Supabase Data Model",
        "12 core tables per user",
        "Postgres + RLS",
        [
            Branch("Identity", ["profile"], TEAL),
            Branch("Career", ["vessels · seatimes · navigation_areas"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Compliance", ["certificates · sea_references"], FOREST),
            Branch("Skills", ["onboard_experiences · specialist_qualifications"], GOLD),
            Branch("Extra", ["tenders · achievements · hobbies_interests"], CORAL),
            Branch("Private", ["payslips"], CORAL),
        ],
        "Every table is user-scoped with row-level security in production.",
        layout="tree",
    ),
    DeckSpec(
        "05-technical-architecture/storage-buckets.md",
        "Storage Buckets",
        "12 private file buckets with signed URLs",
        "Supabase Storage",
        [
            Branch("Photos", ["profile · vessel · tender · hobbies"], TEAL),
            Branch("Documents", ["certificates · references · seatime"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Private files", ["payslips · vessel docs", "Signed URL access", "Never public by default"], CORAL),
            Branch("Upload flow", ["seav-upload.js", "Client-side upload", "RLS enforced"], FOREST),
            Branch("Hardening", ["Public read policies", "Being tightened", "See docs/hardening"], GOLD),
        ],
        "Files stay private until explicitly shown on an opt-in public profile.",
        layout="tree",
    ),
    DeckSpec(
        "05-technical-architecture/auth-and-rls.md",
        "Auth & Security",
        "Login, sessions, and row-level security",
        "Auth + RLS",
        [
            Branch("Auth", ["Email / password", "PKCE sessions", "Forgot password"], TEAL),
            Branch("Protected pages", ["15 app modules", "Redirect if logged out", "auth.js gate"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("RLS", ["user_id on every row", "Phase 2 production", "Anon public read only"], FOREST),
            Branch("Public read", ["public-profile anon client", "Token verify page", "No session leak"], GOLD),
            Branch("Account delete", ["delete_own_account RPC", "Profile page action", "GDPR-ready"], CORAL),
        ],
        "Each user sees only their data. Public viewers see only what crew opt to share.",
    ),
    DeckSpec(
        "05-technical-architecture/deployment-vercel.md",
        "Deployment",
        "How SEA-V ships to production",
        "Vercel",
        [
            Branch("Static site", ["No build step", "HTML + JS served as-is", "Git push deploy"], TEAL),
            Branch("Domain", ["www.sea-v.com", "Vercel project", "Cache headers"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("CSP", ["vercel.json headers", "Supabase + Leaflet allowed", "Security policy"], FOREST),
            Branch("Cache bust", ["ASSET_VERSION in config", "?v= on scripts", "Badge version separate"], GOLD),
            Branch("Analytics", ["Vercel Web Analytics", "Injected in core.js", "Usage insight"], CORAL),
        ],
        "Push to main → live in minutes. No servers to manage.",
    ),
    # --- 06 brand and copy ---
    DeckSpec(
        "06-brand-and-copy/positioning-one-liner.md",
        "Positioning",
        "How to describe SEA-V in one breath",
        "One-liner",
        [
            Branch("What", ["Maritime career platform", "For yacht crew", "Digital profile"], TEAL),
            Branch("Problem", ["Scattered paperwork", "Outdated PDF CVs", "Fragmented sea time"], CORAL),
            Branch("Solution", ["One secure platform", "Track · verify · share", "Modern yacht CV"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Audience", ["Deck · eng · interior", "All career stages", "Life at sea"], FOREST),
        ],
        "SEA-V: the modern digital hub where seafarers manage and present their careers.",
    ),
    DeckSpec(
        "06-brand-and-copy/about-page-copy.md",
        "About Page Copy",
        "Marketing story on about.html",
        "About SEA-V",
        [
            Branch("The idea", ["Modern career management", "Built for seafarers", "Not generic LinkedIn"], TEAL),
            Branch("What it does", ["Sea time · certs · vessels", "References · profile", "Structured records"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Better profile", ["Not a PDF dump", "Transparent progression", "Easy to verify"], GOLD),
            Branch("Vision", ["Just getting started", "Trusted ecosystem", "Industry connections"], FOREST),
        ],
        "Copy tone: professional, crew-first, ambitious but honest about early stage.",
    ),
    DeckSpec(
        "06-brand-and-copy/contact-and-support.md",
        "Contact & Support",
        "How users and partners reach SEA-V",
        "Contact",
        [
            Branch("Email", ["admin@sea-v.com", "Support requests", "Partnership enquiries"], TEAL),
            Branch("Early access", ["Growth phase messaging", "Feedback welcome", "Feature requests"], GOLD),
            Branch("In-app", ["Feedback modal in core.js", "Toast notifications", "User-reported bugs"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Partners", ["Training providers", "Agencies & recruiters", "Industry collab"], FOREST),
        ],
        "Keep the door open — early users and partners shape the roadmap.",
    ),
    DeckSpec(
        "06-brand-and-copy/legal-privacy-terms.md",
        "Legal & Privacy",
        "Trust framework for user data",
        "Legal",
        [
            Branch("Privacy", ["privacy.html", "GDPR-style rights", "Data collection explained"], TEAL),
            Branch("Terms", ["terms.html", "Account rules", "Verification disclaimer"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Signup", ["Legal acceptance checkbox", "Required on register", "Consent captured"], FOREST),
            Branch("Data rights", ["Delete account", "Public profile control", "Export own records"], GOLD),
            Branch("Liability", ["Verification not guarantee", "User-owned content", "See SECURITY.md"], CORAL),
        ],
        "Legal pages build trust. Crew need to know who sees their data and why.",
    ),
    # --- 07 roadmap (continued) ---
    DeckSpec(
        "07-roadmap/next-build.md",
        "Next Build",
        "Priority work after current launch",
        "Next up",
        [
            Branch("Security", ["Storage hardening", "Public read policies", "RLS smoke tests"], CORAL),
            Branch("Public profile", ["Per-section publish", "Anon read path fix", "Empty vs error states"], TEAL),
            Branch("Auth UX", ["No flash before redirect", "Session expiry handling", "Clear on sign-out"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Quality", ["npm lint gate", "Form validation toasts", "User-scoped localStorage"], FOREST),
            Branch("Milestones", ["Vessel attach on auto-award", "Dedupe engine rows", "Badge cache sync"], GOLD),
        ],
        "Harden trust and polish UX before scaling user acquisition.",
    ),
    DeckSpec(
        "07-roadmap/future-networking-verification.md",
        "Future: Network & Verify",
        "Long-range product direction",
        "Future SEA-V",
        [
            Branch("Verification", ["Identity checks", "Sea-time attestation", "Industry stamp of trust"], GOLD),
            Branch("Networking", ["Crew connections", "Department communities", "Mentorship paths"], TEAL),
            Branch("Recruiters", ["Search verified profiles", "Consent-based access", "Agency dashboards"], RGBColor(0x2B, 0x6C, 0xB0)),
            Branch("Integrations", ["Training provider sync", "Certificate auto-import", "Agency placement"], FOREST),
            Branch("Monetisation", ["Pro tier", "Recruiter seats", "Partner revenue share"], CORAL),
        ],
        "From personal career vault to the trusted network layer of the yacht industry.",
    ),
]


def rl_hex(rgb: RGBColor):
    from reportlab.lib.colors import HexColor

    return HexColor(f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}")


def pdf_draw_title_page(c, spec: DeckSpec, w, h):
    from reportlab.lib.colors import white

    c.setFillColor(rl_hex(NAVY))
    c.rect(0, 0, w, h, fill=1, stroke=0)
    c.setFillColor(rl_hex(TEAL))
    c.setFont("Helvetica-Bold", 14)
    c.drawString(48, h - 48, "SEA-V")
    c.setFillColor(white)
    c.setFont("Helvetica-Bold", 32)
    c.drawString(48, h - 160, spec.title)
    c.setFillColor(rl_hex(SKY))
    c.setFont("Helvetica", 16)
    c.drawString(48, h - 200, spec.subtitle)
    c.setFillColor(rl_hex(RGBColor(0x99, 0xAA, 0xBB)))
    c.setFont("Helvetica", 10)
    c.drawString(48, 36, "Website Plan  ·  www.sea-v.com")


def pdf_draw_branch(c, x, y, w, h, branch: Branch, line_color: RGBColor):
    from reportlab.lib.colors import white

    c.setStrokeColor(rl_hex(line_color))
    c.setLineWidth(1.5)
    c.setFillColor(white)
    c.roundRect(x, y, w, h, 10, fill=1, stroke=1)
    c.setFillColor(rl_hex(line_color))
    c.setFont("Helvetica-Bold", 11)
    c.drawString(x + 10, y + h - 18, branch.title[:28])
    c.setFillColor(rl_hex(SLATE))
    c.setFont("Helvetica", 8)
    ty = y + h - 32
    for bullet in branch.bullets[:3]:
        c.drawString(x + 10, ty, f"• {bullet[:38]}")
        ty -= 12


def pdf_draw_radial(c, spec: DeckSpec, w, h):
    from reportlab.lib.colors import white

    c.setFillColor(rl_hex(SKY))
    c.rect(0, 0, w, h, fill=1, stroke=0)
    c.setFillColor(rl_hex(NAVY))
    c.setFont("Helvetica-Bold", 18)
    c.drawString(36, h - 36, spec.title)

    cx, cy = w / 2, h / 2 - 10
    hw, hh = 95, 38
    c.setFillColor(rl_hex(NAVY))
    c.setStrokeColor(rl_hex(TEAL))
    c.setLineWidth(1.5)
    c.roundRect(cx - hw, cy - hh, hw * 2, hh * 2, 12, fill=1, stroke=1)
    c.setFillColor(white)
    c.setFont("Helvetica-Bold", 11)
    c.drawCentredString(cx, cy - 4, spec.center)

    n = len(spec.branches)
    card_w, card_h = 155, 88
    for i, branch in enumerate(spec.branches):
        angle = (2 * math.pi * i / n) - math.pi / 2
        bx = cx + 210 * math.cos(angle) - card_w / 2
        by = cy + 145 * math.sin(angle) - card_h / 2
        color = branch.color or BRANCH_COLORS[i % len(BRANCH_COLORS)]
        tx, ty = bx + card_w / 2, by + card_h / 2
        c.setStrokeColor(rl_hex(color))
        c.setLineWidth(2)
        c.line(cx, cy, tx, ty)
        pdf_draw_branch(c, bx, by, card_w, card_h, branch, color)


def pdf_draw_tree(c, spec: DeckSpec, w, h):
    c.setFillColor(rl_hex(SKY))
    c.rect(0, 0, w, h, fill=1, stroke=0)
    c.setFillColor(rl_hex(NAVY))
    c.setFont("Helvetica-Bold", 18)
    c.drawString(36, h - 36, spec.title)

    root_w, root_h = 180, 34
    rx, ry = w / 2 - root_w / 2, h - 90
    c.setFillColor(rl_hex(NAVY))
    c.setStrokeColor(rl_hex(TEAL))
    c.roundRect(rx, ry, root_w, root_h, 8, fill=1, stroke=1)
    c.setFillColor(rl_hex(RGBColor(0xFF, 0xFF, 0xFF)))
    c.setFont("Helvetica-Bold", 11)
    c.drawCentredString(w / 2, ry + 12, spec.center)

    cols = min(3, len(spec.branches))
    card_w, card_h = 155, 88
    gap_x, gap_y = 175, 105
    start_x = 40
    start_y = h - 220
    for i, branch in enumerate(spec.branches):
        col = i % cols
        row = i // cols
        bx = start_x + col * gap_x
        by = start_y - row * gap_y
        color = branch.color or BRANCH_COLORS[i % len(BRANCH_COLORS)]
        c.setStrokeColor(rl_hex(color))
        c.setLineWidth(1.5)
        c.line(w / 2, ry, bx + card_w / 2, by + card_h)
        pdf_draw_branch(c, bx, by, card_w, card_h, branch, color)


def pdf_draw_takeaway(c, spec: DeckSpec, w, h):
    from reportlab.lib.colors import white

    if not spec.takeaway:
        return
    c.setFillColor(rl_hex(NAVY))
    c.rect(0, 0, w, h, fill=1, stroke=0)
    c.setFillColor(rl_hex(TEAL))
    c.setFont("Helvetica-Bold", 12)
    c.drawString(48, h - 48, "Key takeaway")
    c.setFillColor(white)
    c.setFont("Helvetica-Bold", 22)
    text = c.beginText(48, h - 120)
    text.setLeading(28)
    words = spec.takeaway.split()
    line = ""
    for word in words:
        trial = (line + " " + word).strip()
        if c.stringWidth(trial, "Helvetica-Bold", 22) < w - 96:
            line = trial
        else:
            text.textLine(line)
            line = word
    if line:
        text.textLine(line)
    c.drawText(text)
    c.setFillColor(rl_hex(SKY))
    c.setFont("Helvetica", 12)
    c.drawString(48, 48, spec.subtitle)


def build_pdf(spec: DeckSpec) -> Path:
    from reportlab.lib.pagesizes import landscape, letter
    from reportlab.pdfgen import canvas as pdf_canvas

    out_dir = ROOT / Path(spec.rel_path).parent
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = ROOT / Path(spec.rel_path).with_suffix(".pdf")
    page = landscape(letter)
    w, h = page

    c = pdf_canvas.Canvas(str(out_path), pagesize=page)
    pdf_draw_title_page(c, spec, w, h)
    c.showPage()
    if spec.layout == "tree":
        pdf_draw_tree(c, spec, w, h)
    else:
        pdf_draw_radial(c, spec, w, h)
    c.showPage()
    pdf_draw_takeaway(c, spec, w, h)
    c.showPage()
    c.save()
    return out_path


def write_markdown_stub(spec: DeckSpec):
    md_path = ROOT / spec.rel_path
    md_path.parent.mkdir(parents=True, exist_ok=True)
    lines = [
        f"# {spec.title}",
        "",
        spec.subtitle,
        "",
        f"**Central topic:** {spec.center}",
        "",
    ]
    for branch in spec.branches:
        lines.append(f"## {branch.title}")
        for bullet in branch.bullets:
            lines.append(f"- {bullet}")
        lines.append("")
    lines.append(f"**Key takeaway:** {spec.takeaway}")
    lines.append("")
    md_path.write_text("\n".join(lines), encoding="utf-8")


def main():
    global ROOT
    root = Path(sys.argv[1]) if len(sys.argv) > 1 else ROOT
    ROOT = root
    root.mkdir(parents=True, exist_ok=True)

    for spec in SPECS:
        write_markdown_stub(spec)
        pptx_path = build_deck(spec)
        pdf_path = build_pdf(spec)
        print(f"OK  {pptx_path.relative_to(root)}")
        print(f"OK  {pdf_path.relative_to(root)}")

    print(f"\nGenerated {len(SPECS)} presentation decks + PDFs in {root}")


if __name__ == "__main__":
    main()
